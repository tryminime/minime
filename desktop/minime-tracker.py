#!/usr/bin/env python3
"""
MiniMe System Activity Tracker — Background Daemon

Detects ALL open GUI applications on GNOME Wayland via:
  1. MiniMe GNOME Shell extension D-Bus (focused window + all windows)
  2. /proc/[pid]/exe + /proc/[pid]/cmdline scan (all running GUI apps)

Tracks app open/close events with accurate durations.  Does NOT send
periodic heartbeats — each app session is a single activity record
with start_time and end_time.

Usage:
    python3 minime-tracker.py [--api-url http://localhost:8000] [--interval 2]
"""

import argparse
import configparser
import glob
import hashlib
import json
import logging
import os
import signal
import subprocess
import sys
import time
import urllib.request
import urllib.error
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

# ── Configuration ────────────────────────────────────────────────────────────

POLL_INTERVAL = 2          # seconds between window polls
SYNC_INTERVAL = 30         # seconds between backend syncs
BATCH_SIZE = 50            # max activities per sync
LOG_FILE = "/tmp/minime-tracker.log"
PID_FILE = "/tmp/minime-tracker.pid"

# ── Dynamic App Discovery ────────────────────────────────────────────────────

# Categories considered "productive" — tracked automatically
TRACKED_CATEGORIES = {
    "development", "ide", "texteditor", "webbrowser", "office",
    "graphics", "2dgraphics", "3dgraphics", "rastergraphics", "vectorgraphics",
    "audiovideo", "audio", "video", "player",
    "science", "math", "education", "engineering",
    "network", "email", "chat", "instantmessaging", "ircclient",
    "filetransfer", "p2p", "remoteaccess",
    "photography", "publishing", "wordprocessor", "spreadsheet",
    "presentation", "database", "projectmanagement",
    "terminalemulator", "filemanager",
}

# Categories that indicate system/non-productivity apps → exclude
EXCLUDED_CATEGORIES = {
    "settings", "desktopsettings", "hardwaresettings",
    "packagemanager", "softwaremanagement",
    "security", "documentation", "monitor",
    "filesystem",  # disk usage tools
}

# Interpreters: these run other apps — check their cmdline for real app
INTERPRETER_EXES = {"python3", "python", "python3.10", "python3.11", "python3.12",
                    "java", "javaw", "node", "ruby", "perl"}

# Executables to always skip — background daemons, helpers, shells
SKIP_EXES = {
    "kworker", "ksoftirqd", "kthread", "migration", "bash", "sh", "zsh",
    "fish", "systemd", "dbus-daemon", "dbus-broker", "dbus",
    "pulseaudio", "pipewire", "wireplumber",
    "Xwayland", "gnome-session", "gnome-shell", "gdm", "mutter",
    "npm", "npx", "uvicorn", "gunicorn", "cargo", "rustc", "gcc", "g++", "make",
    "gjs", "gsd-", "gnome-settings-daemon", "at-spi", "at-spi2",
    "polkit", "polkitd", "ibus-daemon", "ibus-engine", "ibus-x11",
    "xdg-desktop-portal", "xdg-document-portal", "xdg-permission-store",
    "gvfsd", "gvfs-udisks2", "tracker-miner", "tracker-extract",
    "fwupd", "snapd", "packagekitd", "update-notifier",
    "chrome_crashpad_handler", "nacl_helper",
    "gnome-software", "seahorse", "yelp", "info",
    "vim", "nvim", "htop", "top",
    "sudo", "pkexec", "env", "cat", "grep", "sed", "awk",
    "snap", "flatpak",
}

# Display names to skip (non-productivity apps that pass category filters)
SKIP_APP_NAMES = {
    "Calendar", "Camera", "Calculator", "Clocks", "Screenshot",
    "Characters", "Fonts", "Help", "Disks", "Logs",
    "Power Statistics", "Startup Applications", "Disk Usage Analyzer",
    "Startup Disk Creator", "Software Updater", "Software & Updates",
    "Input Method", "IBus Preferences", "Language Support",
    "Advanced Network Configuration", "NVIDIA X Server Settings",
    "ImageMagick (color depth=q16)", "Hardware Locality lstopo",
    "Passwords and Keys", "Document Scanner", "Firmware Updater",
    "App Center", "BleachBit", "BleachBit as Administrator",
    "OpenJDK Java 8 Policy Tool",
}

# Manual overrides — exe_name → display_name
# Critical for apps whose binary name differs from .desktop Exec
MANUAL_OVERRIDES = {
    # LibreOffice: runs as soffice.bin but .desktop says libreoffice
    "soffice": "LibreOffice", "soffice.bin": "LibreOffice",
    "oosplash": "LibreOffice",  # LibreOffice splash loader
    "lowriter": "LibreOffice Writer", "localc": "LibreOffice Calc",
    "loimpress": "LibreOffice Impress", "lodraw": "LibreOffice Draw",
    "lomath": "LibreOffice Math", "lobase": "LibreOffice Base",
    # Document viewers
    "evince": "Document Viewer", "evince-previewer": "Document Viewer",
    "okular": "Okular", "xreader": "Document Reader", "atril": "Document Reader",
    "zathura": "Zathura",
    # Browsers
    "brave-browser-stable": "Brave Browser", "brave-browser": "Brave Browser",
    "google-chrome-stable": "Google Chrome", "google-chrome": "Google Chrome",
    "chromium-browser": "Chromium",
    # Electron/Desktop apps
    "electron": "Electron App",
    "obsidian": "Obsidian", "notion-app": "Notion",
    "figma-linux": "Figma",
    "teams": "Microsoft Teams", "zoom": "Zoom",
    "slack": "Slack", "discord": "Discord",
    "spotify": "Spotify", "telegram-desktop": "Telegram",
    # IDEs with different binary names
    "pycharm": "PyCharm", "idea": "IntelliJ IDEA",
    "webstorm": "WebStorm", "clion": "CLion",
    "code": "Visual Studio Code",
    # Text editors
    "gnome-text-editor": "Text Editor", "gedit": "Text Editor",
    "xed": "Text Editor", "pluma": "Text Editor", "mousepad": "Text Editor",
    "kate": "Kate", "kwrite": "KWrite",
}

# ── Master File Category Map ─────────────────────────────────────────────
# Single source of truth for ALL file extension → category mappings.
# Every other extension set in this file is derived from this map.

FILE_CATEGORY_MAP: dict[str, str] = {
    # ── Documents ────────────────────────────────────
    ".pdf": "document", ".doc": "document", ".docx": "document",
    ".odt": "document", ".rtf": "document", ".pages": "document",
    ".epub": "document", ".mobi": "document", ".djvu": "document",

    # ── Notes & Writing ──────────────────────────────
    ".txt": "note", ".md": "note", ".rst": "note",
    ".org": "note", ".tex": "note", ".bib": "note",
    ".adoc": "note", ".wiki": "note",

    # ── Spreadsheets ─────────────────────────────────
    ".xlsx": "spreadsheet", ".xls": "spreadsheet", ".ods": "spreadsheet",
    ".csv": "spreadsheet", ".tsv": "spreadsheet", ".numbers": "spreadsheet",

    # ── Presentations ────────────────────────────────
    ".pptx": "presentation", ".ppt": "presentation", ".odp": "presentation",
    ".key": "presentation",

    # ── Code — General ───────────────────────────────
    ".py": "code", ".js": "code", ".ts": "code", ".jsx": "code",
    ".tsx": "code", ".java": "code", ".kt": "code", ".kts": "code",
    ".cpp": "code", ".c": "code", ".h": "code", ".hpp": "code",
    ".cs": "code", ".swift": "code", ".m": "code",
    ".rs": "code", ".go": "code", ".rb": "code", ".php": "code",
    ".scala": "code", ".clj": "code", ".ex": "code", ".exs": "code",
    ".r": "code", ".R": "code", ".jl": "code", ".lua": "code",
    ".dart": "code", ".zig": "code", ".nim": "code", ".v": "code",
    ".pl": "code", ".pm": "code",

    # ── Code — Web & Markup ──────────────────────────
    ".html": "code", ".htm": "code", ".css": "code", ".scss": "code",
    ".sass": "code", ".less": "code", ".vue": "code", ".svelte": "code",

    # ── Code — Shell & Scripts ───────────────────────
    ".sh": "script", ".bash": "script", ".zsh": "script",
    ".fish": "script", ".ps1": "script", ".bat": "script",
    ".cmd": "script",

    # ── Config & Data ────────────────────────────────
    ".json": "data", ".xml": "data", ".yaml": "data", ".yml": "data",
    ".toml": "data", ".ini": "data", ".cfg": "data", ".conf": "data",
    ".env": "data", ".properties": "data",

    # ── Database & Query ─────────────────────────────
    ".sql": "database", ".sqlite": "database", ".db": "database",

    # ── Notebooks ────────────────────────────────────
    ".ipynb": "notebook", ".rmd": "notebook", ".qmd": "notebook",

    # ── Logs ─────────────────────────────────────────
    ".log": "log",

    # ── Design & Creative ────────────────────────────
    ".psd": "design", ".ai": "design", ".sketch": "design",
    ".fig": "design", ".xd": "design",
    ".blend": "3d_model", ".fbx": "3d_model", ".obj": "3d_model",
    ".xcf": "design", ".odg": "design",

    # ── Images ───────────────────────────────────────
    ".png": "image", ".jpg": "image", ".jpeg": "image",
    ".gif": "image", ".svg": "image", ".webp": "image",
    ".bmp": "image", ".tiff": "image", ".ico": "image",

    # ── Audio ────────────────────────────────────────
    ".mp3": "audio", ".wav": "audio", ".flac": "audio",
    ".ogg": "audio", ".m4a": "audio", ".aac": "audio",

    # ── Video ────────────────────────────────────────
    ".mp4": "video", ".mkv": "video", ".avi": "video",
    ".mov": "video", ".webm": "video", ".flv": "video",

    # ── Archives ─────────────────────────────────────
    ".zip": "archive", ".tar": "archive", ".gz": "archive",
    ".rar": "archive", ".7z": "archive", ".bz2": "archive",
}

# ── Derived Sets (auto-computed from FILE_CATEGORY_MAP) ──────────────────

# All trackable extensions (replaces old FILE_EXTENSIONS)
TRACKABLE_EXTENSIONS = set(FILE_CATEGORY_MAP.keys())

# Reading extensions = document + note + spreadsheet + presentation
READING_EXTENSIONS = {
    ext for ext, cat in FILE_CATEGORY_MAP.items()
    if cat in ("document", "note", "spreadsheet", "presentation")
}

# Categories whose content can be extracted as text for NLP
EXTRACTABLE_CATEGORIES = {
    "document", "note", "spreadsheet", "presentation",
    "code", "script", "data", "database", "notebook", "log",
}

# ── Reading Detection ────────────────────────────────────────────────────

# Known document / reading applications (lowercase match on exe_name or display name)
READING_APPS = {
    # PDF viewers
    "evince", "okular", "zathura", "mupdf", "xpdf", "foxit", "foxitreader",
    "acroread", "qpdfview", "atril", "xreader", "document viewer",
    # Office / document editors
    "libreoffice", "soffice", "soffice.bin", "abiword", "wps", "onlyoffice",
    "libreoffice writer", "libreoffice calc", "libreoffice impress",
    # eBook readers
    "calibre", "foliate", "fbreader", "bookworm", "coolreader",
    # Text editors (count as reading when they have doc files open)
    "gedit", "kate", "mousepad", "pluma", "xed", "leafpad",
    "gnome-text-editor", "text editor",
}


def build_app_map() -> dict[str, str]:
    """Build exe_name → display_name mapping from .desktop files.

    Scans XDG application directories for .desktop files.
    Includes apps whose categories overlap with TRACKED_CATEGORIES.
    Excludes system/non-productivity apps.
    """
    desktop_dirs = [
        "/usr/share/applications",
        os.path.expanduser("~/.local/share/applications"),
        "/var/lib/snapd/desktop/applications",
        "/var/lib/flatpak/exports/share/applications",
        os.path.expanduser("~/.local/share/flatpak/exports/share/applications"),
    ]

    app_map: dict[str, str] = {}

    for d in desktop_dirs:
        if not os.path.isdir(d):
            continue
        for fpath in glob.glob(os.path.join(d, "*.desktop")):
            try:
                cp = configparser.ConfigParser(interpolation=None)
                cp.read(fpath, encoding="utf-8")
                section = "Desktop Entry"
                if not cp.has_section(section):
                    continue
                if cp.get(section, "Type", fallback="") != "Application":
                    continue
                if cp.get(section, "NoDisplay", fallback="false").lower() == "true":
                    continue

                name = cp.get(section, "Name", fallback="").strip()
                exec_line = cp.get(section, "Exec", fallback="").strip()
                categories_raw = cp.get(section, "Categories", fallback="")

                if not name or not exec_line:
                    continue

                # Skip non-productivity apps by name
                if name in SKIP_APP_NAMES:
                    continue

                # Extract binary name from Exec line
                exe_part = exec_line.split()[0]
                exe_name = os.path.basename(exe_part).lower()
                # Handle wrapper prefixes
                for wrapper in ("flatpak", "snap", "env"):
                    if exe_name == wrapper and len(exec_line.split()) > 1:
                        exe_name = os.path.basename(exec_line.split()[1]).lower()

                if not exe_name or exe_name in SKIP_EXES:
                    continue

                # Category filtering
                cats = {c.strip().lower() for c in categories_raw.split(";") if c.strip()}

                # Skip if categories are ONLY system/excluded
                filler = {"gtk", "gnome", "kde", "qt", "x-gnome-utilities",
                          "core", "utility", "x-gnome-personalsettings",
                          "x-gnome-systemsettings", "x-gnome-networksettings",
                          "x-suse-core-office", "x-red-hat-base"}
                meaningful_cats = cats - filler
                if meaningful_cats and meaningful_cats.issubset(EXCLUDED_CATEGORIES):
                    continue

                # Include if any tracked category matches
                if cats & TRACKED_CATEGORIES:
                    app_map[exe_name] = name
                elif not cats:
                    app_map[exe_name] = name

            except Exception:
                continue

    # Add manual overrides
    app_map.update(MANUAL_OVERRIDES)

    return app_map


# ── Logging ──────────────────────────────────────────────────────────────────

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.FileHandler(LOG_FILE),
        logging.StreamHandler(sys.stdout),
    ],
)
log = logging.getLogger("minime-tracker")

# ── Window Detection ─────────────────────────────────────────────────────────

def get_focused_window_dbus() -> Optional[dict]:
    """Get focused window via MiniMe GNOME Shell extension D-Bus."""
    try:
        result = subprocess.run(
            ["gdbus", "call", "--session",
             "--dest", "org.minime.Tracker",
             "--object-path", "/org/minime/Tracker",
             "--method", "org.minime.Tracker.GetFocusedWindow"],
            capture_output=True, text=True, timeout=3,
        )
        if result.returncode != 0:
            return None
        raw = result.stdout.strip()
        json_str = raw.strip("()',\n ")
        data = json.loads(json_str)
        if data.get("app") or data.get("title"):
            return data
    except Exception as e:
        log.debug(f"D-Bus call failed: {e}")
    return None


def get_all_windows_dbus() -> list[dict]:
    """Get ALL open windows via MiniMe GNOME Shell extension D-Bus."""
    try:
        result = subprocess.run(
            ["gdbus", "call", "--session",
             "--dest", "org.minime.Tracker",
             "--object-path", "/org/minime/Tracker",
             "--method", "org.minime.Tracker.GetAllWindows"],
            capture_output=True, text=True, timeout=3,
        )
        if result.returncode != 0:
            return []
        raw = result.stdout.strip()
        json_str = raw.strip("()',\n ")
        return json.loads(json_str)
    except Exception as e:
        log.debug(f"D-Bus GetAllWindows failed: {e}")
    return []


def _extract_file_context(pid: str) -> list[str]:
    """Extract ALL document files a process has open.

    Returns list of absolute file paths.
    Uses /proc/[pid]/fd (best — always absolute) and /proc/[pid]/cmdline.
    """
    found: list[str] = []
    seen: set[str] = set()

    # 1. /proc/fd — gives absolute paths, most reliable
    try:
        fd_dir = f"/proc/{pid}/fd"
        for fd in os.listdir(fd_dir):
            try:
                target = os.readlink(os.path.join(fd_dir, fd))
                if not target.startswith("/"):
                    continue
                if target.startswith(("/dev", "/proc", "/sys", "/run", "/tmp",
                                      "/usr", "/snap", "/var", "/etc")):
                    continue
                ext = Path(target).suffix.lower()
                if ext in TRACKABLE_EXTENSIONS and target not in seen:
                    found.append(target)
                    seen.add(target)
            except (FileNotFoundError, PermissionError, OSError):
                continue
    except (FileNotFoundError, PermissionError):
        pass

    # 2. /proc/cmdline — split on null bytes to preserve filenames with spaces
    try:
        raw = Path(f"/proc/{pid}/cmdline").read_bytes()
        parts = [p.decode("utf-8", errors="replace") for p in raw.split(b"\x00") if p]
        # Get CWD for resolving relative paths
        try:
            cwd = os.readlink(f"/proc/{pid}/cwd")
        except (FileNotFoundError, PermissionError, OSError):
            cwd = None

        for part in parts[1:]:  # Skip the binary itself
            part = part.strip()
            if not part or part.startswith("-"):
                continue  # Skip flags
            # Check if it looks like a file
            ext_match = False
            if "." in part:
                ext = "." + part.rsplit(".", 1)[-1].lower()
                ext_match = ext in TRACKABLE_EXTENSIONS

            if ext_match:
                # Resolve to absolute path
                if part.startswith("/"):
                    full_path = part
                elif cwd:
                    full_path = os.path.join(cwd, part)
                else:
                    full_path = part

                if full_path not in seen:
                    found.append(full_path)
                    seen.add(full_path)
    except (FileNotFoundError, PermissionError):
        pass

    return found


# ── File hash ────────────────────────────────────────────────────────────────

def _compute_file_hash(file_path: str) -> Optional[str]:
    """SHA256 hash of the first 1MB of a file for deduplication."""
    try:
        h = hashlib.sha256()
        with open(file_path, "rb") as f:
            h.update(f.read(1_048_576))
        return h.hexdigest()
    except Exception:
        return None


# ── Content Extraction ───────────────────────────────────────────────────────

EXTRACT_MAX_CHARS = 8_000  # increased from 800 for meaningful NLP


def _extract_content_summary(file_path: str, max_chars: int = EXTRACT_MAX_CHARS) -> Optional[str]:
    """Extract text from a file for NLP analysis.

    Supports PDF, plain text/code, ODT/ODS/ODP, DOCX, XLSX, PPTX.
    """
    if not file_path or not os.path.isfile(file_path):
        return None

    ext = Path(file_path).suffix.lower()
    cat = FILE_CATEGORY_MAP.get(ext)

    if not cat or cat not in EXTRACTABLE_CATEGORIES:
        return None

    try:
        # PDF — uses external tools
        if ext == ".pdf":
            return _extract_pdf_text(file_path, max_chars)

        # OpenDocument formats
        if ext in (".odt", ".ods", ".odp"):
            return _extract_odt_text(file_path, max_chars)

        # Microsoft Word
        if ext in (".docx", ".doc"):
            return _extract_docx_text(file_path, max_chars)

        # Microsoft Excel
        if ext in (".xlsx", ".xls"):
            return _extract_xlsx_text(file_path, max_chars)

        # Microsoft PowerPoint
        if ext in (".pptx", ".ppt"):
            return _extract_pptx_text(file_path, max_chars)

        # Everything else in extractable categories: read as plain text
        # (code, scripts, notes, data/config, database, notebook, log)
        text = Path(file_path).read_text(errors="replace")[:max_chars]
        return text.strip() if text.strip() else None

    except Exception as e:
        log.debug(f"Content extraction failed for {file_path}: {e}")

    return None


def _extract_pdf_text(file_path: str, max_chars: int = EXTRACT_MAX_CHARS) -> Optional[str]:
    """Extract text from PDF — tries pdftotext (all pages), then pdfplumber, then raw."""
    # Try pdftotext first (most reliable, from poppler-utils)
    try:
        result = subprocess.run(
            ["pdftotext", file_path, "-"],  # all pages
            capture_output=True, text=True, timeout=15,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout[:max_chars].strip()
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    # Try pdfplumber (Python library)
    try:
        import pdfplumber  # type: ignore
        texts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages[:20]:  # first 20 pages
                t = page.extract_text()
                if t:
                    texts.append(t)
                if sum(len(x) for x in texts) >= max_chars:
                    break
        combined = "\n".join(texts)[:max_chars]
        if combined.strip():
            return combined
    except Exception:
        pass

    # Fallback: raw PDF stream extraction
    try:
        import re
        with open(file_path, "rb") as f:
            raw = f.read(1_000_000)
        texts = re.findall(rb"\(([^\)]{3,300})\)", raw)
        decoded = []
        for t in texts:
            try:
                s = t.decode("utf-8", errors="replace").strip()
                if s and all(c.isprintable() or c in "\n\t " for c in s) and len(s) > 3:
                    decoded.append(s)
            except Exception:
                continue
        if decoded:
            combined = " ".join(decoded)[:max_chars]
            return combined if len(combined) > 20 else None
    except Exception:
        pass

    # Last resort: PDF metadata
    try:
        import re
        with open(file_path, "rb") as f:
            head = f.read(4096).decode("utf-8", errors="replace")
        parts = []
        for key in ("Title", "Author", "Subject"):
            m = re.search(rf"/{key}\s*\(([^)]+)\)", head)
            if m:
                parts.append(f"{key}: {m.group(1)}")
        if parts:
            return " | ".join(parts)
    except Exception:
        pass

    return None


def _extract_odt_text(file_path: str, max_chars: int = EXTRACT_MAX_CHARS) -> Optional[str]:
    """Extract text from OpenDocument files."""
    try:
        import zipfile, re
        with zipfile.ZipFile(file_path) as z:
            content = z.read("content.xml").decode("utf-8", errors="replace")
            text = re.sub(r"<[^>]+>", " ", content)
            text = re.sub(r"\s+", " ", text).strip()
            return text[:max_chars] if text else None
    except Exception:
        pass
    return None


def _extract_docx_text(file_path: str, max_chars: int = EXTRACT_MAX_CHARS) -> Optional[str]:
    """Extract text from Word .docx files using python-docx."""
    try:
        from docx import Document  # type: ignore
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)[:max_chars]
        return text if text.strip() else None
    except ImportError:
        log.debug("python-docx not installed, trying zipfile fallback for .docx")
    except Exception as e:
        log.debug(f"docx extraction failed: {e}")

    # Fallback: docx is a zip, extract document.xml
    try:
        import zipfile, re
        with zipfile.ZipFile(file_path) as z:
            xml = z.read("word/document.xml").decode("utf-8", errors="replace")
            text = re.sub(r"<[^>]+>", " ", xml)
            text = re.sub(r"\s+", " ", text).strip()
            return text[:max_chars] if text else None
    except Exception:
        pass
    return None


def _extract_xlsx_text(file_path: str, max_chars: int = EXTRACT_MAX_CHARS) -> Optional[str]:
    """Extract cell values from Excel .xlsx files using openpyxl."""
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets[:5]:  # first 5 sheets
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    parts.append(row_text)
                if sum(len(p) for p in parts) >= max_chars:
                    break
            if sum(len(p) for p in parts) >= max_chars:
                break
        text = "\n".join(parts)[:max_chars]
        return text if text.strip() else None
    except ImportError:
        log.debug("openpyxl not installed")
    except Exception as e:
        log.debug(f"xlsx extraction failed: {e}")
    return None


def _extract_pptx_text(file_path: str, max_chars: int = EXTRACT_MAX_CHARS) -> Optional[str]:
    """Extract text from PowerPoint .pptx files using python-pptx."""
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(file_path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            if sum(len(p) for p in parts) >= max_chars:
                break
        text = "\n".join(parts)[:max_chars]
        return text if text.strip() else None
    except ImportError:
        log.debug("python-pptx not installed")
    except Exception as e:
        log.debug(f"pptx extraction failed: {e}")
    return None


def _resolve_interpreter_app(pid: str, app_map: dict[str, str]) -> Optional[str]:
    """For interpreter processes (python3, java, node), check cmdline
    to find the actual application name.

    Example: python3 running spyder → cmdline contains 'spyder'
    """
    try:
        cmdline = Path(f"/proc/{pid}/cmdline").read_text()
        parts = [p for p in cmdline.split("\x00") if p.strip()]
        # Check each cmdline argument against the app map
        for part in parts:
            basename = os.path.basename(part).lower()
            # Direct match
            if basename in app_map:
                return app_map[basename]
            # Partial match (e.g., "spyder" in "/usr/bin/spyder")
            for key, name in app_map.items():
                if key in basename:
                    return name
    except (FileNotFoundError, PermissionError):
        pass
    return None


def get_all_running_apps_proc(app_map: dict[str, str]) -> dict[str, dict]:
    """Scan /proc to find ALL running GUI applications.

    Returns dict mapping app display name -> {app, title, pid, exe, rss, context}.
    Uses /proc/[pid]/exe (+ cmdline for interpreters) to identify apps.
    """
    apps: dict[str, dict] = {}

    try:
        for entry in os.listdir("/proc"):
            if not entry.isdigit():
                continue
            pid = entry

            # Read the exe symlink
            try:
                exe_path = os.readlink(f"/proc/{pid}/exe")
                exe_name = os.path.basename(exe_path).lower()
            except (FileNotFoundError, PermissionError, OSError):
                continue

            # Check if it's an interpreter running a GUI app
            if exe_name in INTERPRETER_EXES:
                resolved = _resolve_interpreter_app(pid, app_map)
                if resolved:
                    # Skip non-productivity
                    if resolved in SKIP_APP_NAMES:
                        continue
                    # Get file context
                    file_paths = _extract_file_context(pid)
                    title = resolved
                    if file_paths:
                        title = f"{resolved} — {os.path.basename(file_paths[0])}"

                    content_summary = None
                    if file_paths:
                        content_summary = _extract_content_summary(file_paths[0])

                    if resolved not in apps:
                        apps[resolved] = {
                            "app": resolved,
                            "title": title,
                            "pid": int(pid),
                            "exe": exe_name,
                            "rss": 0,
                            "file_paths": file_paths,
                            "content_summary": content_summary,
                        }
                    # Update RSS
                    try:
                        stat = Path(f"/proc/{pid}/stat").read_text()
                        rss = int(stat.split()[23])
                        if rss > apps[resolved].get("rss", 0):
                            apps[resolved]["rss"] = rss
                            apps[resolved]["pid"] = int(pid)
                    except (FileNotFoundError, PermissionError, ValueError, IndexError):
                        pass
                continue  # Don't process interpreters further

            # Skip known system/helper processes
            if any(exe_name.startswith(s) for s in SKIP_EXES):
                continue

            # Match against the app map
            display_name = None

            # Exact match
            if exe_name in app_map:
                display_name = app_map[exe_name]
            else:
                # Partial match (e.g., "soffice.bin" → "soffice")
                for key, name in app_map.items():
                    if key in exe_name or exe_name in key:
                        display_name = name
                        break

            if not display_name:
                continue

            # Skip non-productivity apps
            if display_name in SKIP_APP_NAMES:
                continue

            # Get RSS
            try:
                stat = Path(f"/proc/{pid}/stat").read_text()
                fields = stat.split()
                if len(fields) <= 23 or fields[2] == "Z":
                    continue
                rss = int(fields[23])
                if rss < 500:
                    continue
            except (FileNotFoundError, PermissionError, ValueError):
                continue

            # Extract file context (list of absolute paths)
            file_paths = _extract_file_context(pid)
            title = display_name
            if file_paths:
                title = f"{display_name} — {os.path.basename(file_paths[0])}"

            # Extract content summary for the primary file
            content_summary = None
            if file_paths:
                content_summary = _extract_content_summary(file_paths[0])

            # Track highest-RSS process per app
            if display_name not in apps or rss > apps[display_name].get("rss", 0):
                apps[display_name] = {
                    "app": display_name,
                    "title": title,
                    "pid": int(pid),
                    "exe": exe_name,
                    "rss": rss,
                    "file_paths": file_paths,
                    "content_summary": content_summary,
                }
    except Exception as e:
        log.debug(f"/proc scan error: {e}")

    return apps


# ── Activity Tracking ────────────────────────────────────────────────────────

def _get_file_type(path: str) -> str:
    """Get file category from the master FILE_CATEGORY_MAP."""
    ext = Path(path).suffix.lower()
    return FILE_CATEGORY_MAP.get(ext, "file")


def _get_working_dir(pid: int) -> Optional[str]:
    """Get the current working directory of a process."""
    try:
        return os.readlink(f"/proc/{pid}/cwd")
    except (FileNotFoundError, PermissionError, OSError):
        return None


class Activity:
    """Represents a single app usage session (open → close)."""
    def __init__(self, app: str, title: str, pid: int = 0, is_focused: bool = False):
        self.id = hashlib.md5(f"{time.time()}-{app}-{title}".encode()).hexdigest()
        self.app = app
        self.title = title
        self.pid = pid
        self.is_focused = is_focused
        self.started_at = datetime.now(timezone.utc)
        self.ended_at: Optional[datetime] = None
        self.files_accessed: list[dict] = []  # [{path, name, type, content_preview}]
        self.working_dir: Optional[str] = None
        self.content_summary: Optional[str] = None  # Main content preview
        self.content_id: Optional[str] = None       # NLP ingested record linkage

    @staticmethod
    def _sanitize_text(text: Optional[str]) -> Optional[str]:
        """Remove null bytes and control characters that PostgreSQL rejects."""
        if not text:
            return text
        # Remove null bytes (\x00) — asyncpg/PostgreSQL cannot store these
        text = text.replace("\x00", "")
        # Remove other control chars except newline/tab/carriage-return
        text = "".join(c for c in text if c in ("\n", "\r", "\t") or (ord(c) >= 32))
        return text.strip() if text.strip() else None

    def add_file(self, path: str, content_preview: Optional[str] = None):
        """Track a file accessed during this session."""
        name = os.path.basename(path)
        # Avoid duplicates
        if any(f["path"] == path for f in self.files_accessed):
            return
        entry = {
            "path": path,
            "name": name,
            "type": _get_file_type(path),
        }
        if content_preview:
            entry["content_preview"] = self._sanitize_text(content_preview[:500])
        self.files_accessed.append(entry)

    @property
    def duration_seconds(self) -> int:
        end = self.ended_at or datetime.now(timezone.utc)
        return max(0, int((end - self.started_at).total_seconds()))

    @property
    def primary_file(self) -> Optional[dict]:
        return self.files_accessed[0] if self.files_accessed else None

    @property
    def is_reading(self) -> bool:
        """Check if this activity looks like document reading."""
        app_lower = self.app.lower()
        title_lower = self.title.lower()

        # Check if the app is a known reading app
        app_match = any(ra in app_lower for ra in READING_APPS)

        # Check if any accessed file has a reading extension
        has_reading_file = any(
            Path(f["path"]).suffix.lower() in READING_EXTENSIONS
            for f in self.files_accessed
        )

        # Check if the window title contains a document extension
        title_has_doc = any(ext in title_lower for ext in READING_EXTENSIONS)

        # Google Docs / Sheets / Slides in browser
        is_gdocs = any(x in title_lower for x in (
            "google docs", "google sheets", "google slides",
        ))

        return (app_match and (title_has_doc or has_reading_file)) or is_gdocs or has_reading_file

    def to_sync_payload(self) -> dict:
        data = {
            "pid": self.pid,
            "device_id": get_device_id(),
            "started_at": self.started_at.isoformat(),
            "ended_at": (self.ended_at or datetime.now(timezone.utc)).isoformat(),
            "is_focused": self.is_focused,
        }
        if self.files_accessed:
            # Sanitize file entries for DB safety
            clean_files = []
            for f in self.files_accessed:
                cf = {"path": f["path"], "name": f["name"], "type": f.get("type", "file")}
                if f.get("content_preview"):
                    cf["content_preview"] = self._sanitize_text(f["content_preview"][:300])
                clean_files.append(cf)
            data["files"] = clean_files
            data["file_path"] = self.files_accessed[0]["path"]
            data["file_name"] = self.files_accessed[0]["name"]
            data["file_type"] = self.files_accessed[0].get("type", "file")
        if self.content_summary:
            data["content_summary"] = self._sanitize_text(self.content_summary[:500])
        if self.content_id:
            data["content_id"] = self.content_id
        if self.working_dir:
            data["working_directory"] = self.working_dir

        # Build a rich title
        title = self.title
        if self.files_accessed and self.app not in self.title:
            title = f"{self.app} — {self.files_accessed[0]['name']}"
        elif self.files_accessed:
            title = f"{self.app} — {self.files_accessed[0]['name']}"

        # Detect reading vs generic app focus
        activity_type = "reading_analytics" if self.is_reading else "app_focus"

        # Add reading context for reading activities
        if activity_type == "reading_analytics":
            reading_ctx = {
                "time_on_page_sec": self.duration_seconds,
            }
            if self.files_accessed:
                primary = self.files_accessed[0]
                reading_ctx["file_path"] = primary["path"]
                reading_ctx["file_name"] = primary["name"]
                reading_ctx["file_type"] = primary.get("type", "document")
                if primary.get("content_preview"):
                    word_count = len(primary["content_preview"].split())
                    reading_ctx["word_count"] = word_count
            data["reading"] = reading_ctx

        return {
            "type": activity_type,
            "source": "desktop",
            "app": self.app,
            "title": title,
            "duration_seconds": self.duration_seconds,
            "data": data,
        }


def get_device_id() -> str:
    """Stable device identifier."""
    try:
        machine_id = Path("/etc/machine-id").read_text().strip()
    except Exception:
        machine_id = "unknown"
    return hashlib.sha256(f"minime-{machine_id}".encode()).hexdigest()[:16]


# ── Backend Sync ─────────────────────────────────────────────────────────────

class BackendSync:
    """Syncs activities to MiniMe backend API."""

    # Primary: written by the backend on every web/desktop login
    CONFIG_TOKEN_FILE = os.path.expanduser("~/.config/minime/token")
    # Legacy fallback location (kept for backward compat)
    LEGACY_TOKEN_FILE = "/tmp/minime-tracker-token"

    def __init__(self, api_url: str, email: str = None, password: str = None):
        self.api_url = api_url.rstrip("/")
        self.token: Optional[str] = None
        self.refresh_token: Optional[str] = None
        self._email = email or os.environ.get("MINIME_EMAIL")
        self._password = password or os.environ.get("MINIME_PASSWORD")
        self._load_token()

    def _load_token(self):
        """Load auth token — tries sources in priority order.

        Priority:
          1. MINIME_AUTH_TOKEN env var (explicit override)
          2. ~/.config/minime/token  (written by web app on login)
          3. /tmp/minime-tracker-token (legacy location)
          4. Explicit --email/--password flags
        """
        # 1. Env var override
        token = os.environ.get("MINIME_AUTH_TOKEN")
        if token:
            self.token = token
            self.refresh_token = None
            log.info("Auth token loaded from MINIME_AUTH_TOKEN env var")
            return

        # 2. Config file — written automatically on every web login
        config_data = self._try_read_token(self.CONFIG_TOKEN_FILE)
        if config_data:
            if isinstance(config_data, dict):
                self.token = config_data.get("access_token")
                self.refresh_token = config_data.get("refresh_token")
            else:
                self.token = config_data
                self.refresh_token = None
            log.info(f"Auth token loaded from {self.CONFIG_TOKEN_FILE} (web login)")
            # Mirror to legacy location for backward compat
            try:
                Path(self.LEGACY_TOKEN_FILE).write_text(self.token)
            except Exception:
                pass
            return

        # 3. Legacy token file
        legacy_data = self._try_read_token(self.LEGACY_TOKEN_FILE)
        if legacy_data:
            if isinstance(legacy_data, dict):
                self.token = legacy_data.get("access_token")
                self.refresh_token = legacy_data.get("refresh_token")
            else:
                self.token = legacy_data
                self.refresh_token = None
            log.info(f"Auth token loaded from {self.LEGACY_TOKEN_FILE} (legacy)")
            return

        # 4. email/password login as last resort
        if self._email and self._password:
            self._auto_login()
            return

        log.warning(
            "No auth token found. Please log in to MiniMe at http://localhost:3000 — "
            "the tracker will automatically pick up your session token."
        )

    def _try_read_token(self, path: str):
        """Try to read a token file, return dict/str None if missing/empty/invalid."""
        try:
            p = Path(path)
            if p.exists():
                text = p.read_text().strip()
                try:
                    data = json.loads(text)
                    if data.get("access_token"):
                        return data
                except json.JSONDecodeError:
                    if text and len(text) > 20:  # Basic JWT sanity check
                        return text
        except Exception:
            pass
        return None

    def _refresh_session(self) -> bool:
        """Use refresh token to get new access token."""
        if not self.refresh_token:
            return self._auto_login()
        url = f"{self.api_url}/api/v1/auth/refresh"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.refresh_token}"
        }
        try:
            req = urllib.request.Request(url, headers=headers, method="POST", data=b"{}")
            with urllib.request.urlopen(req, timeout=10) as resp:
                body = json.loads(resp.read().decode())
                self.token = body.get("access_token")
                self.refresh_token = body.get("refresh_token", self.refresh_token)
                if self.token:
                    # Save to config dir
                    config_dir = os.path.expanduser("~/.config/minime")
                    os.makedirs(config_dir, exist_ok=True)
                    data = {"access_token": self.token, "refresh_token": self.refresh_token}
                    Path(os.path.join(config_dir, "token")).write_text(json.dumps(data))
                    log.info("Session refreshed successfully")
                    return True
        except Exception as e:
            log.error(f"Token refresh failed: {e}")
            self.refresh_token = None
        return self._auto_login()

    def _auto_login(self) -> bool:
        """Login to backend API and persist the token."""
        if not self._email or not self._password:
            return False
        url = f"{self.api_url}/api/v1/auth/login"
        payload = json.dumps({"email": self._email, "password": self._password}).encode()
        headers = {"Content-Type": "application/json"}
        try:
            req = urllib.request.Request(url, data=payload, headers=headers, method="POST")
            with urllib.request.urlopen(req, timeout=10) as resp:
                body = json.loads(resp.read().decode())
                self.token = body.get("access_token")
                self.refresh_token = body.get("refresh_token")
                if self.token:
                    # Save to config dir so future runs pick it up
                    config_dir = os.path.expanduser("~/.config/minime")
                    os.makedirs(config_dir, exist_ok=True)
                    data = {"access_token": self.token, "refresh_token": self.refresh_token}
                    Path(os.path.join(config_dir, "token")).write_text(json.dumps(data))
                    log.info(f"Logged in as {self._email} — token persisted")
                    return True
        except urllib.error.HTTPError as e:
            body_text = e.read().decode() if e.fp else ""
            log.error(f"Login failed (HTTP {e.code}): {body_text[:200]}")
        except Exception as e:
            log.error(f"Login failed: {e}")
        return False

    def ingest_file(self, file_path: str, text: str, app: str, file_hash: str) -> Optional[str]:
        """Send extracted file text to the backend NLP pipeline.

        Returns the content_id if ingestion succeeded, else None.
        Called from a background thread — must not block the main loop.
        """
        if not self.token:
            return None
        ext = Path(file_path).suffix.lower()
        doc_type = FILE_CATEGORY_MAP.get(ext, "local_file")
        url = f"{self.api_url}/api/v1/content/ingest"
        payload = json.dumps({
            "url": f"file://{file_path}",
            "title": os.path.basename(file_path),
            "doc_type": doc_type,
            "full_text": text,
            "metadata": {
                "file_path": file_path,
                "file_hash": file_hash,
                "app": app,
                "source": "desktop_tracker",
            },
        }).encode("utf-8")
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.token}",
        }
        try:
            req = urllib.request.Request(url, data=payload, headers=headers, method="POST")
            with urllib.request.urlopen(req, timeout=30) as resp:
                body = json.loads(resp.read().decode())
                content_id = body.get("id")
                log.info(f"📚 Content ingested: {os.path.basename(file_path)} → {content_id}")
                return content_id
        except urllib.error.HTTPError as e:
            body_text = e.read().decode() if e.fp else ""
            # 200 with {skipped: true} is fine (localhost filter, duplicate, etc)
            if e.code == 200:
                return None
            log.debug(f"Content ingest HTTP {e.code}: {body_text[:200]}")
            if e.code == 401:
                log.warning("Content ingest 401: Auth token expired — triggering refresh")
                self._refresh_session()
        except Exception as e:
            log.debug(f"Content ingest failed for {file_path}: {e}")
        return None

    def sync(self, activities: list) -> bool:
        """Sync a batch of activities to the backend."""
        if not activities:
            return True
        if not self.token:
            self._load_token()
            if not self.token:
                log.warning("Still no auth token — skipping sync")
                return False

        url = f"{self.api_url}/api/v1/activities/sync"
        payload = {"activities": [a.to_sync_payload() for a in activities]}
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.token}",
        }
        try:
            data = json.dumps(payload).encode("utf-8")
            req = urllib.request.Request(url, data=data, headers=headers, method="POST")
            with urllib.request.urlopen(req, timeout=10) as resp:
                body = json.loads(resp.read().decode())
                log.info(
                    f"Synced {len(activities)} activities → "
                    f"{body.get('synced', '?')} accepted"
                )
                return True
        except urllib.error.HTTPError as e:
            body = e.read().decode() if e.fp else ""
            log.error(f"Sync HTTP error {e.code}: {body[:200]}")
            if e.code == 401:
                log.warning("Sync 401: Auth token expired — refreshing")
                self.token = None
                self._refresh_session()
        except Exception as e:
            log.error(f"Sync failed: {e}")
        return False


# ── Main Loop ────────────────────────────────────────────────────────────────

# Global shutdown flag — signal handler sets this, main loop checks it
_shutdown = False


def write_pid():
    with open(PID_FILE, "w") as f:
        f.write(str(os.getpid()))


def _signal_handler(signum=None, frame=None):
    """Set shutdown flag so main loop exits cleanly with final sync."""
    global _shutdown
    _shutdown = True


def main():
    global _shutdown
    parser = argparse.ArgumentParser(description="MiniMe System Activity Tracker")
    parser.add_argument("--api-url", default="http://localhost:8000", help="Backend API URL")
    parser.add_argument("--interval", type=int, default=POLL_INTERVAL, help="Poll interval (seconds)")
    parser.add_argument("--sync-interval", type=int, default=SYNC_INTERVAL, help="Sync interval (seconds)")
    parser.add_argument("--email", default=None, help="MiniMe account email")
    parser.add_argument("--password", default=None, help="MiniMe account password")
    args = parser.parse_args()

    write_pid()
    signal.signal(signal.SIGTERM, _signal_handler)
    signal.signal(signal.SIGINT, _signal_handler)

    log.info(f"MiniMe tracker starting (API: {args.api_url}, poll: {args.interval}s, sync: {args.sync_interval}s)")

    # Check D-Bus extension
    ext_available = get_focused_window_dbus() is not None
    if ext_available:
        log.info("✓ GNOME extension detected — using D-Bus for window tracking")
    else:
        log.warning("⚠ GNOME extension not loaded — using /proc scan for ALL running apps")

    # Build exe→name mapping from .desktop files
    app_map = build_app_map()
    log.info(f"📦 Discovered {len(app_map)} trackable apps from .desktop files")

    sync_client = BackendSync(args.api_url, email=args.email, password=args.password)

    # State: track all currently open apps
    # Key = app display name, Value = Activity object (with start time)
    active_apps: dict[str, Activity] = {}
    pending_sync: list[Activity] = []  # Closed activities waiting to be synced
    last_sync_time = time.time()

    # ── Initial detection ────────────────────────────────────────────
    if ext_available:
        windows = get_all_windows_dbus()
        for win in windows:
            app = win.get("app", "")
            if app and app not in SKIP_APP_NAMES and app not in active_apps:
                title = win.get("title", app)
                active_apps[app] = Activity(app=app, title=title,
                                            pid=win.get("pid", 0),
                                            is_focused=win.get("focus", False))
    else:
        proc_apps = get_all_running_apps_proc(app_map)
        for name, info in proc_apps.items():
            act = Activity(app=name, title=info.get("title", name),
                           pid=info["pid"], is_focused=False)
            # Add all discovered files with content previews
            for fp in info.get("file_paths", []):
                preview = _extract_content_summary(fp)
                act.add_file(fp, content_preview=preview)
            act.content_summary = info.get("content_summary")
            act.working_dir = _get_working_dir(info["pid"])
            active_apps[name] = act

    if active_apps:
        names = ", ".join(sorted(active_apps.keys()))
        log.info(f"🖥️  {len(active_apps)} apps detected: {names}")
    else:
        log.warning("No apps detected at startup")

    # ── Main polling loop ────────────────────────────────────────────
    while not _shutdown:
        try:
            # Detect current apps
            current_apps: dict[str, dict] = {}

            if ext_available:
                windows = get_all_windows_dbus()
                for win in windows:
                    app = win.get("app", "")
                    if app and app not in SKIP_APP_NAMES:
                        current_apps[app] = {
                            "title": win.get("title", app),
                            "pid": win.get("pid", 0),
                            "focus": win.get("focus", False),
                            "context": None,
                        }
            else:
                proc_apps = get_all_running_apps_proc(app_map)
                for name, info in proc_apps.items():
                    current_apps[name] = {
                        "title": info.get("title", name),
                        "pid": info["pid"],
                        "focus": False,
                        "file_paths": info.get("file_paths", []),
                        "content_summary": info.get("content_summary"),
                    }

            # ── Detect newly opened apps ─────────────────────────────
            for app_name, info in current_apps.items():
                if app_name not in active_apps:
                    act = Activity(
                        app=app_name,
                        title=info["title"],
                        pid=info["pid"],
                        is_focused=info.get("focus", False),
                    )
                    # Capture file context
                    file_paths = info.get("file_paths", [])
                    for fp in file_paths:
                        text = _extract_content_summary(fp)
                        act.add_file(fp, content_preview=text)
                        # Ingest into NLP pipeline in background thread
                        if text and len(text) > 100:
                            fhash = _compute_file_hash(fp)
                            if fhash:
                                import threading
                                def _do_ingest(path=fp, t=text, a=app_name, h=fhash, act_ref=act):
                                    cid = sync_client.ingest_file(path, t, a, h)
                                    if cid:
                                        act_ref.content_id = cid
                                threading.Thread(target=_do_ingest, daemon=True).start()
                    act.content_summary = info.get("content_summary")
                    # Capture working directory
                    act.working_dir = _get_working_dir(info["pid"])
                    active_apps[app_name] = act
                    file_info = f" [{os.path.basename(file_paths[0])}]" if file_paths else ""
                    log.info(f"🟢 App opened: {app_name}{file_info}")
                else:
                    # Update context if new files appeared
                    existing = active_apps[app_name]
                    for fp in info.get("file_paths", []):
                        if fp not in [f.get("path") for f in existing.files_accessed]:
                            text = _extract_content_summary(fp)
                            existing.add_file(fp, content_preview=text)
                            # Ingest new file in background
                            if text and len(text) > 100:
                                fhash = _compute_file_hash(fp)
                                if fhash:
                                    import threading
                                    def _do_ingest_update(path=fp, t=text, a=app_name, h=fhash, act_ref=existing):
                                        cid = sync_client.ingest_file(path, t, a, h)
                                        if cid and not act_ref.content_id:
                                            act_ref.content_id = cid
                                    threading.Thread(target=_do_ingest_update, daemon=True).start()
                    # Update title if changed
                    new_title = info.get("title", app_name)
                    if new_title != existing.title:
                        existing.title = new_title

            # ── Detect closed apps ───────────────────────────────────
            closed = [name for name in active_apps if name not in current_apps]
            for app_name in closed:
                activity = active_apps.pop(app_name)
                if activity.duration_seconds >= 3:
                    activity.ended_at = datetime.now(timezone.utc)
                    pending_sync.append(activity)
                    log.info(
                        f"🔴 App closed: {app_name} "
                        f"(used for {activity.duration_seconds}s)"
                    )

            # ── Periodic sync of closed activities ───────────────────
            now = time.time()
            if now - last_sync_time >= args.sync_interval:
                # Sync completed activities
                if pending_sync:
                    batch = pending_sync[:BATCH_SIZE]
                    if sync_client.sync(batch):
                        pending_sync = pending_sync[BATCH_SIZE:]

                # Also sync snapshots of ACTIVE activities so dashboard shows
                # real-time data (not just after an app closes)
                active_snapshots = []
                for app_name, act in active_apps.items():
                    if act.duration_seconds >= 5:  # Only report if used > 5s
                        active_snapshots.append(act)
                if active_snapshots:
                    sync_client.sync(active_snapshots)

                last_sync_time = now

            time.sleep(args.interval)

        except KeyboardInterrupt:
            _shutdown = True
        except Exception as e:
            log.error(f"Tracker error: {e}")
            time.sleep(5)

    # ── Final sync on exit ───────────────────────────────────────────
    # Close all active apps and sync them
    for app_name, activity in active_apps.items():
        if activity.duration_seconds >= 3:
            activity.ended_at = datetime.now(timezone.utc)
            pending_sync.append(activity)
            log.info(f"📋 Closing: {app_name} ({activity.duration_seconds}s)")

    if pending_sync:
        log.info(f"Final sync: {len(pending_sync)} activities")
        sync_client.sync(pending_sync)

    # Cleanup PID file
    try:
        os.unlink(PID_FILE)
    except FileNotFoundError:
        pass
    log.info("MiniMe tracker stopped")


if __name__ == "__main__":
    main()
